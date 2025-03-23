from flask import Flask, render_template, request, send_file
import os
from pptx import Presentation
import fitz  # PyMuPDF para manejar PDFs

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def extraer_texto_pptx(pptx_path):
    """Extrae texto de un archivo PPTX."""
    prs = Presentation(pptx_path)
    texto_extraido = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto_extraido.append(shape.text)
    return "\n".join(texto_extraido)

def extraer_texto_pdf(pdf_path):
    """Extrae texto de un archivo PDF."""
    doc = fitz.open(pdf_path)
    texto_extraido = []
    for page in doc:
        texto_extraido.append(page.get_text())
    return "\n".join(texto_extraido)

@app.route("/", methods=["GET", "POST"])
def index():
    texto = None
    if request.method == "POST":
        if "file" not in request.files:
            return "No se subió ningún archivo"

        file = request.files["file"]
        if file.filename == "":
            return "Selecciona un archivo"

        filepath = os.path.join(app.config["UPLOAD_FOLDER"], file.filename)
        file.save(filepath)
        
        if file.filename.endswith(".pptx"):
            texto = extraer_texto_pptx(filepath)
        elif file.filename.endswith(".pdf"):
            texto = extraer_texto_pdf(filepath)
        else:
            return "Formato no soportado. Solo se permiten archivos PPTX y PDF."

        with open("extraido.txt", "w", encoding="utf-8") as f:
            f.write(texto)
    
    return render_template("index.html", texto=texto)

@app.route("/descargar")
def descargar():
    return send_file("extraido.txt", as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)